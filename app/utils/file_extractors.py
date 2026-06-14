from pathlib import Path

from docx import Document
from pptx import Presentation

from app.services.local_image_extractor_service import extract_image_with_LOCAL
from app.services.llm_router_service import (
    AI_MODE_TEST_LOCAL_ONLY,
    PROVIDER_SKIP,
    TASK_VISION_EXTRACT,
    resolve_provider_for_task,
)
from app.services.portal_ai_mode_service import get_current_portal_ai_mode

IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".gif",
}

IMAGE_VISION_SKIPPED_MESSAGE = (
    "Image vision analysis skipped because local vision is not available "
    "for the selected AI mode."
)


def extract_txt_text(
    file_path: Path
) -> str:

    return file_path.read_text(
        encoding="utf-8"
    )


def extract_docx_text(
    file_path: Path
) -> str:

    doc = Document(file_path)

    return "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
        if paragraph.text.strip()
    )
    

def _is_image_file(file_path: str | Path) -> bool:
    return Path(file_path).suffix.lower() in IMAGE_EXTENSIONS



def _current_ai_mode() -> str:
    portal_ai_mode = get_current_portal_ai_mode()

    if portal_ai_mode and portal_ai_mode.get("ai_mode"):
        return str(portal_ai_mode["ai_mode"]).strip().upper()

    import os

    return (
        os.getenv("NON_PORTAL_AI_MODE")
        or os.getenv("TELEGRAM_AI_MODE")
        or os.getenv("PORTAL_DEFAULT_AI_MODE")
        or "NO_LLM"
    ).strip().upper()


def _extract_image_text(file_path: str | Path) -> str:
    ai_mode = _current_ai_mode()
    provider = resolve_provider_for_task(
        TASK_VISION_EXTRACT,
        ai_mode,
    )

    if provider.get("provider") == PROVIDER_SKIP:
        if ai_mode == AI_MODE_TEST_LOCAL_ONLY:
            raise RuntimeError(
                "Local Vision is required in TEST_LOCAL_ONLY mode but is "
                "not available. Set LOCAL_BASE_URL and LOCAL_VISION_MODEL "
                "for your Ollama/Qwen-VL server."
            )

        return IMAGE_VISION_SKIPPED_MESSAGE

    return extract_image_with_LOCAL(file_path)


def extract_pptx_text(
    file_path: Path
) -> str:

    prs = Presentation(
        file_path
    )

    output = []

    image_index = 1

    image_output_dir = (
        file_path.parent
        / "extracted_images"
    )

    image_output_dir.mkdir(
        parents=True,
        exist_ok=True
    )

    for slide_index, slide in enumerate(
        prs.slides,
        start=1
    ):

        slide_parts = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_parts.append(
                        text
                    )

            if hasattr(shape, "image"):

                image = shape.image

                image_path = (
                    image_output_dir
                    / f"slide_{slide_index}_image_{image_index}.{image.ext}"
                )

                image_path.write_bytes(
                    image.blob
                )

                image_text = _extract_image_text(
                    image_path
                )

                if image_text:
                    slide_parts.append(
                        "# Extracted Image Text\n"
                        + image_text
                    )

                image_index += 1

        if slide_parts:
            output.append(
                f"# Slide {slide_index}\n"
                + "\n\n".join(slide_parts)
            )

    return "\n\n".join(output)


def extract_file_text(
    file_path: Path
) -> str:

    suffix = file_path.suffix.lower()

    if suffix in [".txt", ".md"]:
        return extract_txt_text(
            file_path
        )

    if suffix == ".docx":
        return extract_docx_text(
            file_path
        )

    if suffix == ".pptx":
        return extract_pptx_text(
            file_path
        )

    if suffix in IMAGE_EXTENSIONS:
        return _extract_image_text(
            file_path
        )

    raise ValueError(
        f"Unsupported file type: {suffix}"
    )
