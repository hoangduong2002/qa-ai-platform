import json
import logging
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Optional

from fastapi import UploadFile

logger = logging.getLogger(__name__)

from app.services.jira_requirement_service import (
    create_requirement_from_jira,
)
from app.services.jira_delta_service import (
    build_and_save_latest_stored_jira_snapshot,
    load_latest_change_impact_report,
)
from app.services.impact_mapping_service import (
    load_latest_regeneration_plan,
)
from app.services.requirement_source_service import has_jira_snapshot
from app.services.portal_job_service import update_job_progress

from app.services.requirement_sanitization_service import (
    sanitize_requirement_for_analysis,
)
from app.utils.clarification_answers import (
    count_matched_clarification_answers,
    get_clarification_answer_text,
    get_clarification_id,
    get_clarification_question,
    merge_clarifications_with_answers,
    normalize_clarification_answers,
    normalize_question_text,
)

from graph.nodes.load_requirement import (
    load_requirement,
)

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils import get_column_letter


REQUIREMENTS_DIR = Path("requirements")

def _read_json(
    file_path: Path,
):
    if not file_path.exists():
        return None

    try:
        return json.loads(
            file_path.read_text(
                encoding="utf-8",
                errors="ignore",
            )
        )
    except Exception:
        return None


def _read_latest_versioned_json(
    directory: Path,
    glob_pattern: str,
    version_pattern: str,
    default,
):
    if not directory.exists():
        return default

    latest_path = None
    latest_version = 0

    for path in directory.glob(glob_pattern):
        match = re.match(version_pattern, path.name)
        if match and int(match.group(1)) > latest_version:
            latest_version = int(match.group(1))
            latest_path = path

    if latest_path is None:
        return default

    value = _read_json(latest_path)
    return default if value is None else value


def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def _safe_requirement_id(value: str) -> str:
    value = value.strip().replace(" ", "-")
    safe = []

    for char in value:
        if char.isalnum() or char in ["-", "_"]:
            safe.append(char)

    result = "".join(safe).upper()

    if not result:
        raise ValueError("Requirement name is invalid.")

    return result


def _requirement_dir(ticket_id: str) -> Path:
    return REQUIREMENTS_DIR / ticket_id


def _source_dir(ticket_id: str) -> Path:
    return _requirement_dir(ticket_id) / "source"


def _analysis_dir(ticket_id: str) -> Path:
    return _requirement_dir(ticket_id) / "analysis"


def _ticket_json_file(ticket_id: str) -> Path:
    return _requirement_dir(ticket_id) / "ticket.json"


def _write_source_metadata(ticket_id: str, values: dict) -> None:
    metadata_file = _requirement_dir(ticket_id) / "metadata.json"
    metadata = _read_json(metadata_file) or {}
    metadata.update(values)
    metadata_file.write_text(
        json.dumps(metadata, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def _mark_web_jira_requirement(ticket_id: str, jira_key: str) -> None:
    now = _now_iso()

    ticket_file = _ticket_json_file(ticket_id)
    ticket_data = _read_json(ticket_file) or {}
    ticket_data.update(
        {
            "ticket_id": ticket_id,
            "source": "jira",
            "source_type": "jira",
            "source_channel": "web",
            "imported_from_jira": True,
            "jira_key": jira_key,
            "updated_at": now,
        }
    )
    ticket_file.write_text(
        json.dumps(ticket_data, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    _write_source_metadata(
        ticket_id,
        {
            "ticket_id": ticket_id,
            "source": f"jira:{jira_key}",
            "source_type": "jira",
            "source_channel": "web",
            "imported_from_jira": True,
            "jira_key": jira_key,
            "updated_at": now,
        },
    )


def list_requirements() -> list[dict]:
    if not REQUIREMENTS_DIR.exists():
        return []

    items = []

    for path in REQUIREMENTS_DIR.iterdir():
        if not path.is_dir():
            continue

        if path.name.startswith("_"):
            continue

        ticket_file = path / "ticket.json"

        if ticket_file.exists():
            try:
                data = json.loads(
                    ticket_file.read_text(encoding="utf-8")
                )
            except Exception:
                data = {}
        else:
            data = {}

        items.append(
            {
                "ticket_id": path.name,
                "summary": data.get("summary", path.name),
                "source": data.get("source", "unknown"),
                "created_at": data.get("created_at", ""),
                "updated_at": data.get("updated_at", ""),
                "has_analysis": (path / "analysis" / "requirement_analysis.json").exists(),
                "has_sanitized": (path / "analysis" / "sanitized_requirement.md").exists(),
            }
        )

    return sorted(
        items,
        key=lambda item: item.get("updated_at") or item.get("created_at") or "",
        reverse=True,
    )


async def create_manual_requirement(
    requirement_name: str,
    description: str,
    files: Optional[List[UploadFile]] = None,
) -> str:
    ticket_id = _safe_requirement_id(requirement_name)

    base_dir = _requirement_dir(ticket_id)

    if base_dir.exists():
        raise ValueError(f"Requirement already exists: {ticket_id}")

    source_dir = _source_dir(ticket_id)
    uploads_dir = source_dir / "uploads"
    extracted_dir = source_dir / "extracted"
    analysis_dir = _analysis_dir(ticket_id)

    source_dir.mkdir(parents=True, exist_ok=True)
    uploads_dir.mkdir(parents=True, exist_ok=True)
    extracted_dir.mkdir(parents=True, exist_ok=True)
    analysis_dir.mkdir(parents=True, exist_ok=True)

    description_file = source_dir / "description.md"
    comments_file = source_dir / "comments.md"

    description_file.write_text(
        description.strip(),
        encoding="utf-8",
    )

    comments_file.write_text(
        "",
        encoding="utf-8",
    )

    uploaded_content_parts = []

    if files:
        for upload_file in files:
            if not upload_file.filename:
                continue

            saved_file = uploads_dir / upload_file.filename

            content = await upload_file.read()

            saved_file.write_bytes(content)

            extracted_text = extract_file_text(saved_file)

            extracted_file = extracted_dir / f"{saved_file.stem}.txt"

            extracted_file.write_text(
                extracted_text,
                encoding="utf-8",
            )

            if extracted_text.strip():
                uploaded_content_parts.append(
                    f"## File: {upload_file.filename}\n\n{extracted_text}"
                )

    uploaded_content = "\n\n".join(uploaded_content_parts)

    if uploaded_content:
        uploaded_content_file = source_dir / "uploaded_content.md"
        uploaded_content_file.write_text(
            uploaded_content,
            encoding="utf-8",
        )

    ticket_data = {
        "ticket_id": ticket_id,
        "summary": requirement_name.strip(),
        "source": "web_manual",
        "source_type": "manual",
        "source_channel": "web",
        "imported_from_jira": False,
        "created_at": _now_iso(),
        "updated_at": _now_iso(),
    }

    _ticket_json_file(ticket_id).write_text(
        json.dumps(ticket_data, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    (base_dir / "metadata.json").write_text(
        json.dumps(
            {
                "ticket_id": ticket_id,
                "source": "web_manual",
                "source_type": "manual",
                "source_channel": "web",
                "imported_from_jira": False,
                "created_at": ticket_data["created_at"],
                "updated_at": ticket_data["updated_at"],
            },
            indent=2,
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    sanitize_existing_requirement(ticket_id)

    return ticket_id


def create_requirement_from_jira_and_sanitize(
    issue_key: str,
    jira_pat: str = "",
    refresh_existing: bool = False,
) -> str:
    issue_key = issue_key.strip()
    ticket_id = _safe_requirement_id(issue_key)

    if requirement_is_complete(ticket_id) and not refresh_existing:
        logger.info(
            "Requirement exists and is complete, skipping Jira load. ticket_id=%s",
            ticket_id,
        )
        _mark_web_jira_requirement(ticket_id, ticket_id)
        if not has_jira_snapshot(ticket_id):
            try:
                build_and_save_latest_stored_jira_snapshot(ticket_id)
            except Exception as error:
                logger.warning(
                    "Could not build Jira snapshot for existing requirement. ticket_id=%s error=%s",
                    ticket_id,
                    error,
                )
        return ticket_id

    if _requirement_dir(ticket_id).exists() and not requirement_is_complete(ticket_id):
        logger.warning(
            "Requirement folder exists but is incomplete, rebuilding. ticket_id=%s",
            ticket_id,
        )
        refresh_existing = True

    ticket_id = create_requirement_from_jira(
        issue_key.strip(),
        jira_pat=jira_pat.strip(),
        refresh_existing=refresh_existing,
        source_channel="web",
    )
    _mark_web_jira_requirement(ticket_id, ticket_id)

    update_job_progress(
        current_step="Sanitizing requirement",
        message="Cleaning the Jira requirement for analysis and AI workflows.",
    )

    sanitize_existing_requirement(ticket_id)

    build_and_save_latest_stored_jira_snapshot(ticket_id)

    return ticket_id


def sanitize_existing_requirement(
    ticket_id: str,
) -> str:
    state = {
        "ticket_id": ticket_id,
    }

    state.update(
        load_requirement(state)
    )

    requirement_context = state.get("requirement_context", "")

    sanitized = sanitize_requirement_for_analysis(
        ticket_id=ticket_id,
        raw_requirement=requirement_context,
    )

    return sanitized


def get_requirement_detail(
    ticket_id: str,
) -> dict:
    base_dir = _requirement_dir(ticket_id)
    source_dir = _source_dir(ticket_id)
    analysis_dir = _analysis_dir(ticket_id)

    ticket_file = _ticket_json_file(ticket_id)
    description_file = source_dir / "description.md"
    comments_file = source_dir / "comments.md"
    uploaded_content_file = source_dir / "uploaded_content.md"
    sanitized_file = analysis_dir / "sanitized_requirement.md"

    requirement_analysis_file = analysis_dir / "requirement_analysis.json"
    requirement_items_file = analysis_dir / "requirement_items.json"
    clarifications_file = analysis_dir / "clarifications.json"
    clarification_snapshot_file = analysis_dir / "clarification_questions_snapshot.json"
    clarification_answers_file = analysis_dir / "clarification_answers.json"
    requirement_summary_file = analysis_dir / "requirement_summary.json"

    if not base_dir.exists():
        raise FileNotFoundError(
            f"Requirement not found: {ticket_id}"
        )

    ticket_data = {}

    if ticket_file.exists():
        ticket_data = json.loads(
            ticket_file.read_text(
                encoding="utf-8",
            )
        )

    requirement_analysis = _read_json(
        requirement_analysis_file
    )

    requirement_items = _read_json(
        requirement_items_file
    )

    clarifications_raw = (
        _read_json(clarification_snapshot_file)
        or _read_json(clarifications_file)
    )

    clarifications = _normalize_clarifications(
        clarifications_raw
    )
    clarification_answers = normalize_clarification_answers(
        _read_json(clarification_answers_file)
    )
    clarifications_with_answers = merge_clarifications_with_answers(
        clarifications,
        clarification_answers,
    )
    answered_clarification_count = sum(
        1
        for item in clarifications_with_answers
        if item.get("answer")
    )

    requirement_summary = _read_json(
        requirement_summary_file
    )

    analysis_error_file = analysis_dir / "analyze_error.txt"
    analysis_error = _read_text(analysis_error_file)
    incremental_requirement_items = _read_latest_versioned_json(
        analysis_dir,
        "incremental_requirement_items_v*.json",
        r"incremental_requirement_items_v(\d+)\.json$",
        [],
    )
    incremental_scenario_merge_report = _read_latest_versioned_json(
        analysis_dir,
        "incremental_scenario_merge_report_v*.json",
        r"incremental_scenario_merge_report_v(\d+)\.json$",
        {},
    )
    incremental_testcase_merge_report = _read_latest_versioned_json(
        analysis_dir,
        "incremental_testcase_merge_report_v*.json",
        r"incremental_testcase_merge_report_v(\d+)\.json$",
        {},
    )
    has_incremental_testcases = (
        base_dir / "generated" / "latest_testcases.json"
    ).exists()

    return {
        "ticket_id": ticket_id,
        "ticket": ticket_data,
        "description": _read_text(description_file),
        "comments": _read_text(comments_file),
        "uploaded_content": _read_text(uploaded_content_file),
        "sanitized_requirement": _read_text(sanitized_file),

        "requirement_analysis": requirement_analysis,
        "requirement_items": requirement_items,
        "clarifications": clarifications_with_answers,
        "clarification_answers": clarification_answers,
        "clarifications_with_answers": clarifications_with_answers,
        "answered_clarification_count": answered_clarification_count,
        "unanswered_clarification_count": max(
            len(clarifications_with_answers) - answered_clarification_count,
            0,
        ),
        "requirement_summary": requirement_summary,
        "change_impact_report": load_latest_change_impact_report(ticket_id),
        "regeneration_plan": load_latest_regeneration_plan(ticket_id),
        "incremental_requirement_items": incremental_requirement_items,
        "incremental_scenario_merge_report": incremental_scenario_merge_report,
        "incremental_testcase_merge_report": incremental_testcase_merge_report,
        "analysis_error": analysis_error,

        "has_sanitized": sanitized_file.exists(),
        "has_analysis": requirement_analysis_file.exists(),
        "has_items": requirement_items_file.exists(),
        "has_summary": requirement_summary_file.exists(),
        "has_clarifications": len(clarifications_with_answers) > 0,
        "has_incremental_testcases": has_incremental_testcases,
    }


def update_requirement(
    ticket_id: str,
    summary: str,
    description: str,
    comments: str = "",
) -> None:
    base_dir = _requirement_dir(ticket_id)
    source_dir = _source_dir(ticket_id)

    if not base_dir.exists():
        raise FileNotFoundError(f"Requirement not found: {ticket_id}")

    (source_dir / "description.md").write_text(
        description.strip(),
        encoding="utf-8",
    )

    (source_dir / "comments.md").write_text(
        comments.strip(),
        encoding="utf-8",
    )

    ticket_file = _ticket_json_file(ticket_id)

    ticket_data = {}
    if ticket_file.exists():
        ticket_data = json.loads(
            ticket_file.read_text(encoding="utf-8")
        )

    ticket_data["summary"] = summary.strip()
    ticket_data["updated_at"] = _now_iso()

    ticket_file.write_text(
        json.dumps(ticket_data, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    sanitize_existing_requirement(ticket_id)


def delete_requirement(
    ticket_id: str,
) -> None:
    base_dir = _requirement_dir(ticket_id)

    if not base_dir.exists():
        raise FileNotFoundError(f"Requirement not found: {ticket_id}")

    deleted_dir = REQUIREMENTS_DIR / "_deleted"
    deleted_dir.mkdir(parents=True, exist_ok=True)

    target_dir = deleted_dir / f"{ticket_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    shutil.move(
        str(base_dir),
        str(target_dir),
    )


def extract_file_text(
    file_path: Path,
) -> str:
    suffix = file_path.suffix.lower()

    if suffix in [".txt", ".md", ".csv", ".json", ".xml", ".log"]:
        return file_path.read_text(
            encoding="utf-8",
            errors="ignore",
        )

    if suffix == ".docx":
        return _extract_docx(file_path)

    if suffix == ".pptx":
        return _extract_pptx(file_path)

    return f"[Unsupported file type for text extraction: {file_path.name}]"


def _extract_docx(
    file_path: Path,
) -> str:
    try:
        from docx import Document
    except ImportError:
        return "[python-docx is not installed. Run: pip install python-docx]"

    document = Document(str(file_path))

    paragraphs = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    return "\n".join(paragraphs)


def _extract_pptx(
    file_path: Path,
) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[python-pptx is not installed. Run: pip install python-pptx]"

    presentation = Presentation(str(file_path))

    texts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts.append(f"# Slide {slide_index}")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                value = shape.text.strip()
                if value:
                    texts.append(value)

    return "\n\n".join(texts)


def _read_text(
    file_path: Path,
) -> str:
    if not file_path.exists():
        return ""

    return file_path.read_text(
        encoding="utf-8",
        errors="ignore",
    )
    

def get_clarification_questions(
    ticket_id: str,
) -> list[dict]:
    analysis_dir = _analysis_dir(ticket_id)

    clarification_file = (
        analysis_dir / "clarifications.json"
    )
    clarification_snapshot_file = (
        analysis_dir / "clarification_questions_snapshot.json"
    )
    clarification_answers_file = (
        analysis_dir / "clarification_answers.json"
    )

    data = _read_json(
        clarification_snapshot_file
    )
    if not data:
        data = _read_json(
            clarification_file
        )

    if not data:
        return []

    questions = []

    if isinstance(data, list):
        questions = [_normalize_clarification_question(item) for item in data]
    elif isinstance(data, dict):
        if isinstance(data.get("clarification_questions"), list):
            questions = [
                _normalize_clarification_question(item)
                for item in data["clarification_questions"]
            ]
        elif isinstance(data.get("questions"), list):
            questions = [
                _normalize_clarification_question(item)
                for item in data["questions"]
            ]
        elif isinstance(data.get("clarifications"), list):
            questions = [
                _normalize_clarification_question(item)
                for item in data["clarifications"]
            ]

    if not questions:
        return []

    answers = normalize_clarification_answers(
        _read_json(clarification_answers_file)
    )

    return merge_clarifications_with_answers(
        questions,
        answers,
    )


def _normalize_clarification_question(item: dict) -> dict:
    if not isinstance(item, dict):
        return {}

    question = dict(item)
    question_id = (
        question.get("question_id")
        or question.get("clarification_id")
        or question.get("id")
        or ""
    )

    if question_id:
        question["id"] = question.get("id") or question_id
        question["question_id"] = question_id

    question["question"] = (
        question.get("question")
        or question.get("question_text")
        or question.get("text")
        or question.get("description")
        or ""
    )
    question["impact"] = question.get("impact") or "Medium"
    question["category"] = question.get("category") or "Other"
    question["reason"] = question.get("reason") or ""
    question["free_text_allowed"] = question.get("free_text_allowed", True)

    if not isinstance(question.get("suggested_options"), list):
        question["suggested_options"] = []

    return question


def save_clarification_answers(
    ticket_id: str,
    answers: dict,
) -> None:
    base_dir = _requirement_dir(ticket_id)
    source_dir = _source_dir(ticket_id)
    analysis_dir = _analysis_dir(ticket_id)

    if not base_dir.exists():
        raise FileNotFoundError(
            f"Requirement not found: {ticket_id}"
        )

    analysis_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    source_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    questions = get_clarification_questions(
        ticket_id
    )

    answer_items = []
    submitted_answer_count = len(answers)
    answered_at = datetime.utcnow().isoformat(timespec="seconds") + "Z"
    used_question_ids = set()

    for question in questions:
        question_id = (
            question.get("question_id")
            or question.get("id")
            or ""
        )
        used_question_ids.add(question_id)

        answer_payload = answers.get(question_id, {})

        if isinstance(answer_payload, str):
            answer_payload = {"answer": answer_payload}

        if not isinstance(answer_payload, dict):
            answer_payload = {}

        selected_option_key = str(
            answer_payload.get("selected_option_key") or ""
        ).strip()
        custom_answer = str(
            answer_payload.get("custom_answer")
            or answer_payload.get("answer")
            or ""
        ).strip()

        selected_option_label = ""

        for option in question.get("suggested_options", []):
            if not isinstance(option, dict):
                continue

            option_key = str(option.get("key") or "").strip()

            if option_key == selected_option_key:
                selected_option_label = str(option.get("label") or "").strip()
                break

        final_answer = custom_answer or selected_option_label

        answer_items.append(
            {
                "id": question.get("id") or question_id,
                "question_id": question_id,
                "question": question.get("question", ""),
                "answer": final_answer,
                "answered_at": answered_at if final_answer else "",
                "category": question.get("category", "Other"),
                "impact": question.get("impact", "Medium"),
                "reason": question.get("reason", ""),
                "priority": question.get("priority", ""),
                "related_requirement": question.get("related_requirement", ""),
                "free_text_allowed": question.get("free_text_allowed", True),
                "selected_option_key": selected_option_key,
                "selected_option_label": selected_option_label,
                "custom_answer": custom_answer,
                "final_answer": final_answer,
            }
        )

    for question_id, answer_payload in answers.items():
        if question_id in used_question_ids:
            continue

        if isinstance(answer_payload, str):
            answer_payload = {"answer": answer_payload}

        if not isinstance(answer_payload, dict):
            answer_payload = {}

        final_answer = str(
            answer_payload.get("custom_answer")
            or answer_payload.get("answer")
            or ""
        ).strip()

        answer_items.append(
            {
                "id": question_id,
                "question_id": question_id,
                "question": "",
                "answer": final_answer,
                "answered_at": answered_at if final_answer else "",
                "selected_option_key": str(
                    answer_payload.get("selected_option_key") or ""
                ).strip(),
                "selected_option_label": "",
                "custom_answer": final_answer,
                "final_answer": final_answer,
            }
        )

    answers_file = (
        analysis_dir / "clarification_answers.json"
    )
    output = {
        "ticket_id": ticket_id,
        "answers": answer_items,
    }

    answers_file.write_text(
        json.dumps(
            output,
            indent=2,
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    notes = _build_clarification_answer_notes(
        answer_items
    )

    notes_file = (
        source_dir / "clarification_answer_notes.md"
    )

    notes_file.write_text(
        notes,
        encoding="utf-8",
    )

    summary_file = analysis_dir / "requirement_summary.json"
    if summary_file.exists():
        summary_file.unlink()

    matched_answer_count = count_matched_clarification_answers(
        questions,
        answer_items,
    )
    logger.debug(
        "Saved clarification answers",
        extra={
            "ticket_id": ticket_id,
            "submitted_answer_count": submitted_answer_count,
            "saved_answer_count": len(answer_items),
            "matched_answer_count": matched_answer_count,
            "saved_answer_file_path": str(answers_file),
        },
    )


def _build_clarification_answer_notes(
    answer_items: list[dict],
) -> str:
    lines = [
        "# Clarification Answer Notes",
        "",
        "The following answers were provided by the user and must be treated as additional requirement context.",
        "",
    ]

    for item in answer_items:
        answer = (
            item.get("final_answer")
            or item.get("answer")
            or ""
        ).strip()

        if not answer:
            continue

        lines.extend(
            [
                f"## {item.get('question_id', '')}",
                "",
                f"Question: {item.get('question', '')}",
                "",
                f"Impact: {item.get('impact', 'N/A')}",
                "",
            ]
        )

        selected_option_key = item.get("selected_option_key", "")

        if selected_option_key:
            lines.extend(
                [
                    f"Selected option: {selected_option_key}",
                    "",
                ]
            )

        lines.extend(
            [
                f"Answer: {answer}",
                "",
            ]
        )

    return "\n".join(lines).strip()


def export_requirement_summary_to_excel(
    ticket_id: str,
) -> Path:
    analysis_dir = _analysis_dir(ticket_id)

    summary_file = analysis_dir / "requirement_summary.json"

    if not summary_file.exists():
        raise FileNotFoundError(
            f"Requirement summary not found for {ticket_id}"
        )

    requirement_summary = _read_json(summary_file) or {}

    export_dir = analysis_dir / "exports"
    export_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    output_file = export_dir / f"{ticket_id}_requirement_summary.xlsx"

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Requirement Summary"

    headers = [
        "Section",
        "ID",
        "Description",
        "Priority",
    ]

    worksheet.append(headers)

    header_fill = PatternFill(
        start_color="1F2937",
        end_color="1F2937",
        fill_type="solid",
    )

    header_font = Font(
        color="FFFFFF",
        bold=True,
    )

    for cell in worksheet[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(
            vertical="center",
            horizontal="center",
        )

    def append_items(
        section_name: str,
        items,
    ):
        if not items:
            return

        if isinstance(items, str):
            worksheet.append(
                [
                    section_name,
                    "",
                    items,
                    "",
                ]
            )
            return

        if isinstance(items, dict):
            worksheet.append(
                [
                    section_name,
                    items.get("id", ""),
                    (
                        items.get("description")
                        or items.get("text")
                        or items.get("requirement")
                        or str(items)
                    ),
                    items.get("priority", ""),
                ]
            )
            return

        if isinstance(items, list):
            for item in items:
                if isinstance(item, dict):
                    worksheet.append(
                        [
                            section_name,
                            item.get("id", ""),
                            (
                                item.get("description")
                                or item.get("text")
                                or item.get("requirement")
                                or item.get("rule")
                                or item.get("validation")
                                or item.get("integration")
                                or item.get("error")
                                or str(item)
                            ),
                            item.get("priority", ""),
                        ]
                    )
                else:
                    worksheet.append(
                        [
                            section_name,
                            "",
                            str(item),
                            "",
                        ]
                    )

    overview = requirement_summary.get("overview")

    if overview:
        append_items(
            "Overview",
            overview,
        )

    append_items(
        "Functional Requirements",
        requirement_summary.get("functional_requirements"),
    )

    append_items(
        "Business Rules",
        requirement_summary.get("business_rules"),
    )

    append_items(
        "Validations",
        requirement_summary.get("validations"),
    )

    append_items(
        "Integrations",
        requirement_summary.get("integrations"),
    )

    append_items(
        "Error Handling",
        requirement_summary.get("error_handling"),
    )

    append_items(
        "Non-functional Requirements",
        requirement_summary.get("non_functional_requirements"),
    )

    append_items(
        "Assumptions",
        requirement_summary.get("assumptions"),
    )

    append_items(
        "Out of Scope",
        requirement_summary.get("out_of_scope"),
    )

    append_items(
        "Open Questions",
        requirement_summary.get("open_questions"),
    )

    column_widths = {
        "A": 28,
        "B": 16,
        "C": 100,
        "D": 16,
    }

    for column, width in column_widths.items():
        worksheet.column_dimensions[column].width = width

    for row in worksheet.iter_rows():
        for cell in row:
            cell.alignment = Alignment(
                vertical="top",
                wrap_text=True,
            )

    worksheet.freeze_panes = "A2"

    workbook.save(output_file)

    return output_file

def export_requirement_analysis_to_excel(
    ticket_id: str,
) -> Path:
    analysis_dir = _analysis_dir(ticket_id)
    source_dir = _source_dir(ticket_id)

    requirement_analysis_file = analysis_dir / "requirement_analysis.json"
    requirement_items_file = analysis_dir / "requirement_items.json"
    clarifications_file = analysis_dir / "clarifications.json"
    clarification_snapshot_file = analysis_dir / "clarification_questions_snapshot.json"
    clarification_answers_file = analysis_dir / "clarification_answers.json"
    requirement_summary_file = analysis_dir / "requirement_summary.json"
    sanitized_requirement_file = analysis_dir / "sanitized_requirement.md"

    if not requirement_analysis_file.exists():
        raise FileNotFoundError(
            f"Requirement analysis not found for {ticket_id}"
        )

    requirement_analysis = _read_json(requirement_analysis_file) or {}
    requirement_items = _read_json(requirement_items_file) or []
    clarifications_raw = (
        _read_json(clarification_snapshot_file)
        or _read_json(clarifications_file)
    )

    clarifications = _normalize_clarifications(
        clarifications_raw
    )
    clarification_answers = normalize_clarification_answers(
        _read_json(clarification_answers_file)
    )
    requirement_summary = _read_json(requirement_summary_file) or {}
    sanitized_requirement = _read_text(sanitized_requirement_file)

    if not sanitized_requirement:
        sanitized_requirement = _read_text(source_dir / "description.md")

    export_dir = analysis_dir / "exports"
    export_dir.mkdir(
        parents=True,
        exist_ok=True,
    )

    output_file = export_dir / f"{ticket_id}_analysis_result.xlsx"

    workbook = Workbook()

    default_sheet = workbook.active
    workbook.remove(default_sheet)

    header_fill = PatternFill(
        start_color="1F2937",
        end_color="1F2937",
        fill_type="solid",
    )

    header_font = Font(
        color="FFFFFF",
        bold=True,
    )

    title_font = Font(
        bold=True,
        size=14,
    )

    def style_header_row(worksheet, row_number: int = 1):
        for cell in worksheet[row_number]:
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(
                vertical="center",
                horizontal="center",
                wrap_text=True,
            )

    def auto_format(worksheet):
        for row in worksheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(
                    vertical="top",
                    wrap_text=True,
                )

        for column_cells in worksheet.columns:
            max_length = 0
            column_letter = get_column_letter(column_cells[0].column)

            for cell in column_cells:
                value = cell.value
                if value is None:
                    continue

                max_length = max(
                    max_length,
                    min(len(str(value)), 80),
                )

            worksheet.column_dimensions[column_letter].width = max(
                min(max_length + 2, 60),
                14,
            )

        worksheet.freeze_panes = "A2"

    def append_key_value_sheet(
        sheet_name: str,
        data: dict,
    ):
        worksheet = workbook.create_sheet(sheet_name)

        worksheet.append(
            [
                "Section",
                "Value",
            ]
        )
        style_header_row(worksheet)

        for key, value in data.items():
            if isinstance(value, list):
                value_text = "\n".join(
                    [
                        str(item)
                        for item in value
                    ]
                )
            elif isinstance(value, dict):
                value_text = json.dumps(
                    value,
                    indent=2,
                    ensure_ascii=False,
                )
            else:
                value_text = str(value)

            worksheet.append(
                [
                    key,
                    value_text,
                ]
            )

        auto_format(worksheet)

    def append_items_sheet(
        sheet_name: str,
        items,
        headers: list[str],
        row_builder,
    ):
        worksheet = workbook.create_sheet(sheet_name)

        worksheet.append(headers)
        style_header_row(worksheet)

        if isinstance(items, dict):
            possible_items = (
                items.get("items")
                or items.get("requirements")
                or items.get("clarifications")
                or items.get("questions")
                or []
            )
        else:
            possible_items = items

        if not possible_items:
            worksheet.append(
                [
                    "No data",
                ]
            )
        else:
            for index, item in enumerate(possible_items, start=1):
                worksheet.append(
                    row_builder(
                        item,
                        index,
                    )
                )

        auto_format(worksheet)

    def normalize_text_item(item):
        if isinstance(item, dict):
            return (
                item.get("description")
                or item.get("text")
                or item.get("requirement")
                or item.get("rule")
                or item.get("validation")
                or item.get("integration")
                or item.get("error")
                or json.dumps(
                    item,
                    ensure_ascii=False,
                )
            )

        return str(item)

    # Sheet 1: Overview
    overview_sheet = workbook.create_sheet("Overview")
    overview_sheet.append(
        [
            "Field",
            "Value",
        ]
    )
    style_header_row(overview_sheet)

    overview_sheet.append(
        [
            "Ticket ID",
            ticket_id,
        ]
    )

    overview_sheet.append(
        [
            "Requirement Items",
            len(requirement_items)
            if isinstance(requirement_items, list)
            else 0,
        ]
    )

    overview_sheet.append(
        [
            "Clarification Questions",
            len(clarifications)
            if isinstance(clarifications, list)
            else 0,
        ]
    )

    overview_sheet.append(
        [
            "Has Requirement Summary",
            "Yes" if requirement_summary else "No",
        ]
    )

    overview_sheet.append(
        [
            "Export Type",
            "Requirement Analysis Result",
        ]
    )

    auto_format(overview_sheet)

    # Sheet 2: Requirement Analysis
    append_key_value_sheet(
        "Requirement Analysis",
        requirement_analysis,
    )

    # Sheet 3: Requirement Items
    append_items_sheet(
        "Requirement Items",
        requirement_items,
        [
            "ID",
            "Type",
            "Requirement",
            "Priority",
        ],
        lambda item, index: [
            item.get("id")
            or item.get("requirement_id")
            or item.get("item_id")
            or f"REQ-{index:03d}"
            if isinstance(item, dict)
            else f"REQ-{index:03d}",
            item.get("type")
            or item.get("category")
            or item.get("requirement_type")
            or "N/A"
            if isinstance(item, dict)
            else "N/A",
            normalize_text_item(item),
            item.get("priority", "N/A")
            if isinstance(item, dict)
            else "N/A",
        ],
    )

    # Sheet 4: Clarifications
    answers_by_id = {}
    answers_by_question = {}

    for item in clarification_answers:
        question_id = get_clarification_id(item)
        question_text = normalize_question_text(get_clarification_question(item))

        if question_id:
            answers_by_id[question_id] = item
        if question_text:
            answers_by_question[question_text] = item

    matched_answer_keys = set()
    matched_answer_indexes = set()
    merged_clarifications = []

    def format_cell_value(value):
        if isinstance(value, (list, dict)):
            return json.dumps(value, ensure_ascii=False)
        return value or ""

    for index, item in enumerate(clarifications, start=1):
        if not isinstance(item, dict):
            merged_clarifications.append(
                {
                    "question_id": f"Q{index:03d}",
                    "question": str(item),
                    "answer": "",
                    "answer_status": "Unanswered",
                }
            )
            continue

        question_id = get_clarification_id(item) or f"Q{index:03d}"
        question_text = get_clarification_question(item)
        answer_info = answers_by_id.get(question_id, {})

        if not answer_info:
            answer_info = answers_by_question.get(
                normalize_question_text(question_text),
                {},
            )

        answer = get_clarification_answer_text(answer_info) if answer_info else ""

        if answer_info:
            matched_answer_indexes.add(id(answer_info))
            answer_id = get_clarification_id(answer_info)
            answer_question = normalize_question_text(
                get_clarification_question(answer_info)
            )
            if answer_id:
                matched_answer_keys.add(("id", answer_id))
            if answer_question:
                matched_answer_keys.add(("question", answer_question))

        merged_clarifications.append(
            {
                "question_id": question_id,
                "question": question_text,
                "answer": answer,
                "answer_status": "Answered" if answer else "Unanswered",
                "answered_at": answer_info.get("answered_at", "") if answer_info else "",
                "impact": item.get("impact") or item.get("reason", ""),
                "priority": item.get("priority") or item.get("severity", ""),
                "related_requirement": format_cell_value(
                    item.get("related_requirement")
                    or item.get("related_requirement_id")
                    or item.get("related_requirement_ids")
                    or item.get("requirement_id")
                    or ""
                ),
                "category": item.get("category") or item.get("type", ""),
            }
        )

    for item in clarification_answers:
        question_id = get_clarification_id(item)
        question_text = get_clarification_question(item)
        normalized_question = normalize_question_text(question_text)

        if (
            question_id
            and ("id", question_id) in matched_answer_keys
        ) or (
            normalized_question
            and ("question", normalized_question) in matched_answer_keys
        ):
            continue

        merged_clarifications.append(
            {
                "question_id": question_id,
                "question": question_text,
                "answer": get_clarification_answer_text(item),
                "answer_status": "Answered",
                "answered_at": item.get("answered_at", ""),
                "impact": item.get("impact") or item.get("reason", ""),
                "priority": item.get("priority") or item.get("severity", ""),
                "related_requirement": format_cell_value(
                    item.get("related_requirement")
                    or item.get("related_requirement_id")
                    or item.get("related_requirement_ids")
                    or item.get("requirement_id")
                    or ""
                ),
                "category": item.get("category") or item.get("type", ""),
            }
        )

    def build_clarification_row(item, index: int):
        if not isinstance(item, dict):
            return [
                f"Q{index:03d}",
                str(item),
                "",
                "Unanswered",
                "",
                "",
                "",
                "",
                "",
            ]

        return [
            item.get("question_id", ""),
            item.get("question", ""),
            item.get("answer", ""),
            item.get("answer_status", ""),
            item.get("answered_at", ""),
            item.get("impact", ""),
            item.get("priority", ""),
            item.get("related_requirement", ""),
            item.get("category", ""),
        ]

    append_items_sheet(
        "Clarifications",
        merged_clarifications,
        [
            "Question ID",
            "Question",
            "Answer",
            "Answer Status",
            "Answered At",
            "Impact / Reason",
            "Priority",
            "Related Requirement",
            "Category",
        ],
        build_clarification_row,
    )

    clarification_export_stats = {
        "clarification_count": len(clarifications),
        "answer_count": len(clarification_answers),
        "matched_answer_count": len(matched_answer_indexes),
    }

    # Sheet 5: Requirement Summary, if available
    if requirement_summary:
        summary_sheet = workbook.create_sheet("Requirement Summary")

        summary_sheet.append(
            [
                "Section",
                "ID",
                "Description",
                "Priority",
            ]
        )
        style_header_row(summary_sheet)

        def append_summary_items(section_name: str, items):
            if not items:
                return

            if isinstance(items, str):
                summary_sheet.append(
                    [
                        section_name,
                        "",
                        items,
                        "",
                    ]
                )
                return

            if isinstance(items, dict):
                summary_sheet.append(
                    [
                        section_name,
                        items.get("id", ""),
                        normalize_text_item(items),
                        items.get("priority", ""),
                    ]
                )
                return

            if isinstance(items, list):
                for item in items:
                    if isinstance(item, dict):
                        summary_sheet.append(
                            [
                                section_name,
                                item.get("id", ""),
                                normalize_text_item(item),
                                item.get("priority", ""),
                            ]
                        )
                    else:
                        summary_sheet.append(
                            [
                                section_name,
                                "",
                                str(item),
                                "",
                            ]
                        )

        append_summary_items(
            "Overview",
            requirement_summary.get("overview"),
        )
        append_summary_items(
            "Functional Requirements",
            requirement_summary.get("functional_requirements"),
        )
        append_summary_items(
            "Business Rules",
            requirement_summary.get("business_rules"),
        )
        append_summary_items(
            "Validations",
            requirement_summary.get("validations"),
        )
        append_summary_items(
            "Integrations",
            requirement_summary.get("integrations"),
        )
        append_summary_items(
            "Error Handling",
            requirement_summary.get("error_handling"),
        )
        append_summary_items(
            "Non-functional Requirements",
            requirement_summary.get("non_functional_requirements"),
        )
        append_summary_items(
            "Assumptions",
            requirement_summary.get("assumptions"),
        )
        append_summary_items(
            "Out of Scope",
            requirement_summary.get("out_of_scope"),
        )
        append_summary_items(
            "Open Questions",
            requirement_summary.get("open_questions"),
        )

        auto_format(summary_sheet)

    # Sheet 6: Sanitized Requirement
    sanitized_sheet = workbook.create_sheet("Sanitized Requirement")
    sanitized_sheet.append(
        [
            "Requirement Context",
        ]
    )
    style_header_row(sanitized_sheet)

    sanitized_sheet.append(
        [
            sanitized_requirement,
        ]
    )

    sanitized_sheet.column_dimensions["A"].width = 120
    sanitized_sheet["A2"].alignment = Alignment(
        vertical="top",
        wrap_text=True,
    )

    workbook.save(output_file)
    logger.debug(
        "Exported Clarifications sheet",
        extra={
            "ticket_id": ticket_id,
            **clarification_export_stats,
            "export_path": str(output_file),
        },
    )

    return output_file


def _normalize_clarifications(
    clarifications,
) -> list:
    if not clarifications:
        return []

    if isinstance(clarifications, list):
        return [
            _normalize_clarification_question(item)
            for item in clarifications
        ]

    if isinstance(clarifications, dict):
        for key in [
            "clarification_questions",
            "questions",
            "clarifications",
            "items",
            "data",
        ]:
            value = clarifications.get(key)

            if isinstance(value, list):
                return [
                    _normalize_clarification_question(item)
                    for item in value
                ]

    return []


def requirement_exists(
    ticket_id: str,
) -> bool:
    ticket_id = _safe_requirement_id(ticket_id)

    return _requirement_dir(ticket_id).exists()


def requirement_is_complete(ticket_id: str) -> bool:
    """Check whether a requirement has been fully loaded and sanitized.

    Minimum criteria for a complete requirement:
      - ticket.json exists
      - source/jira_requirement.md or source/description.md exists
      - analysis/sanitized_requirement.md exists
    """
    ticket_id = _safe_requirement_id(ticket_id)
    base_dir = _requirement_dir(ticket_id)

    if not base_dir.exists():
        return False

    ticket_file = _ticket_json_file(ticket_id)
    if not ticket_file.exists():
        return False

    source_dir = _source_dir(ticket_id)
    has_jira_md = (source_dir / "jira_requirement.md").exists()
    has_description_md = (source_dir / "description.md").exists()
    if not has_jira_md and not has_description_md:
        return False

    analysis_dir = _analysis_dir(ticket_id)
    if not (analysis_dir / "sanitized_requirement.md").exists():
        return False

    return True


def normalize_requirement_id(
    value: str,
) -> str:
    return _safe_requirement_id(value)
